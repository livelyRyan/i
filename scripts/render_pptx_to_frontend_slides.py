#!/usr/bin/env python3
"""
Render a PPTX into a zero-build, self-contained HTML slide deck (frontend-slides style),
using absolute positioning to keep layout close to the original PPT.

Notes:
- This does NOT rasterize the full slide background (PowerPoint shapes/gradients). It extracts:
  - text boxes (including those inside groups)
  - pictures (including those inside groups)
  - auto shapes / lines as rough rectangles/lines when possible
- Pictures are clickable (lightbox) for full-size viewing.
"""

from __future__ import annotations

import html
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400
PX_PER_INCH = 96
PT_TO_PX = PX_PER_INCH / 72.0


def emu_to_px(emu: int) -> float:
    return (emu / EMU_PER_INCH) * PX_PER_INCH


@dataclass(frozen=True)
class Transform:
    sx: float = 1.0
    sy: float = 1.0
    tx: float = 0.0
    ty: float = 0.0

    def map_point(self, x: float, y: float) -> Tuple[float, float]:
        return (self.tx + self.sx * x, self.ty + self.sy * y)

    def map_size(self, w: float, h: float) -> Tuple[float, float]:
        return (self.sx * w, self.sy * h)


@dataclass
class Item:
    kind: str  # text | image | shape | line
    x: float
    y: float
    w: float
    h: float
    rotation: float = 0.0
    z: int = 0
    text: str = ""
    font_px: Optional[float] = None
    bold: bool = False
    color: Optional[str] = None
    bg: Optional[str] = None
    border: Optional[str] = None
    border_w: Optional[float] = None
    src: Optional[str] = None
    alt: str = ""


def _rgb_to_hex(rgb) -> str:
    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"


def _extract_text(shape) -> Tuple[str, Optional[float], bool, Optional[str]]:
    """
    Returns: (text, max_font_px, any_bold, color_hex)
    """
    if not getattr(shape, "has_text_frame", False):
        return ("", None, False, None)

    raw = shape.text or ""
    raw = raw.strip()
    if not raw:
        return ("", None, False, None)

    max_pt = 0.0
    any_bold = False
    color_hex = None

    try:
        tf = shape.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                if r.font.bold:
                    any_bold = True
                if r.font.size is not None:
                    max_pt = max(max_pt, float(r.font.size.pt))
                if color_hex is None and r.font.color is not None and r.font.color.type == MSO_COLOR_TYPE.RGB:
                    color_hex = _rgb_to_hex(r.font.color.rgb)
    except Exception:
        pass

    max_px = (max_pt * PT_TO_PX) if max_pt > 0 else None
    # Normalize line breaks for HTML.
    safe = html.escape(raw).replace("\n", "<br/>")
    return (safe, max_px, any_bold, color_hex)


def _extract_shape_style(shape) -> Tuple[Optional[str], Optional[str], Optional[float]]:
    """
    Returns: (bg_color, border_color, border_width_px)
    """
    bg = None
    border = None
    border_w = None

    try:
        if shape.fill is not None and shape.fill.type == MSO_FILL_TYPE.SOLID:
            fc = shape.fill.fore_color
            if fc.type == MSO_COLOR_TYPE.RGB:
                bg = _rgb_to_hex(fc.rgb)
    except Exception:
        pass

    try:
        if shape.line is not None and shape.line.color is not None:
            lc = shape.line.color
            if lc.type == MSO_COLOR_TYPE.RGB:
                border = _rgb_to_hex(lc.rgb)
        if shape.line is not None and shape.line.width is not None:
            border_w = float(shape.line.width) * PT_TO_PX / 12700.0  # EMU-ish; best-effort
            # The above is a heuristic; keep it readable.
            border_w = max(1.0, min(4.0, border_w))
    except Exception:
        pass

    return bg, border, border_w


def _group_child_transform(group_shape, parent: Transform) -> Transform:
    """
    Compute transform mapping from group's child coordinate system to slide EMU coordinates,
    then apply parent transform (for nested groups).
    """
    xfrm = group_shape._element.xpath(".//a:xfrm")[0]
    off = xfrm.off
    ext = xfrm.ext
    ch_off = xfrm.chOff
    ch_ext = xfrm.chExt

    # Derived from:
    # x_slide = parent.tx + parent.sx*(off.x + (x_child - ch_off.x)/ch_ext.cx * ext.cx)
    sx = parent.sx * (ext.cx / ch_ext.cx)
    sy = parent.sy * (ext.cy / ch_ext.cy)
    tx = parent.tx + parent.sx * off.x - sx * ch_off.x
    ty = parent.ty + parent.sy * off.y - sy * ch_off.y
    return Transform(sx=sx, sy=sy, tx=tx, ty=ty)


def _flatten_shapes(shapes: Iterable, t: Transform) -> List[Tuple[object, Transform]]:
    """
    Returns a flat list of (shape, transform_for_shape_local_coords_to_slide_emu).
    For normal shapes, the transform is the current t.
    For children of group shapes, the transform is derived from group geometry.
    """
    out: List[Tuple[object, Transform]] = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            child_t = _group_child_transform(shape, t)
            out.extend(_flatten_shapes(shape.shapes, child_t))
        else:
            out.append((shape, t))
    return out


def _shape_box(shape, t: Transform) -> Tuple[float, float, float, float]:
    x, y = t.map_point(float(shape.left), float(shape.top))
    w, h = t.map_size(float(shape.width), float(shape.height))
    return x, y, w, h


def _write_image(shape, slide_idx: int, pic_idx: int, assets_dir: Path) -> str:
    image = shape.image
    ext = (image.ext or "png").lower()
    filename = f"slide{slide_idx:02d}_pic{pic_idx:02d}.{ext}"
    path = assets_dir / filename
    path.write_bytes(image.blob)
    return filename


def build_deck(pptx_path: Path, out_html: Path, assets_dir: Path) -> None:
    prs = Presentation(str(pptx_path))
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    base_w_px = emu_to_px(slide_w)
    base_h_px = emu_to_px(slide_h)
    ratio = base_w_px / base_h_px if base_h_px else (16 / 9)

    assets_dir.mkdir(parents=True, exist_ok=True)

    slides_items: List[List[Item]] = []
    all_images: List[str] = []

    for s_i, slide in enumerate(prs.slides, start=1):
        flat = _flatten_shapes(slide.shapes, Transform())
        items: List[Item] = []
        pic_idx = 0
        z = 0

        for shape, t in flat:
            z += 1
            rot = float(getattr(shape, "rotation", 0.0) or 0.0)
            x, y, w, h = _shape_box(shape, t)

            # Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic_idx += 1
                filename = _write_image(shape, s_i, pic_idx, assets_dir)
                all_images.append(filename)
                items.append(
                    Item(
                        kind="image",
                        x=x,
                        y=y,
                        w=w,
                        h=h,
                        rotation=rot,
                        z=z,
                        src=f"{assets_dir.name}/{filename}",
                        alt=getattr(shape, "name", f"slide{s_i} image {pic_idx}"),
                    )
                )
                continue

            # Lines
            if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                _, border, border_w = _extract_shape_style(shape)
                items.append(
                    Item(
                        kind="line",
                        x=x,
                        y=y,
                        w=w,
                        h=h,
                        rotation=rot,
                        z=z,
                        border=border or "#ffffff",
                        border_w=border_w or 2.0,
                    )
                )
                continue

            # Text (text boxes, auto shapes with text, etc.)
            text, font_px, any_bold, color_hex = _extract_text(shape)
            bg, border, border_w = _extract_shape_style(shape) if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE else (None, None, None)

            if text:
                items.append(
                    Item(
                        kind="text",
                        x=x,
                        y=y,
                        w=w,
                        h=h,
                        rotation=rot,
                        z=z,
                        text=text,
                        font_px=font_px,
                        bold=any_bold,
                        color=color_hex,
                        bg=bg,
                        border=border,
                        border_w=border_w,
                    )
                )
            else:
                # Auto-shapes without text: best-effort rectangle render to keep layout closer.
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    bg2, border2, border_w2 = _extract_shape_style(shape)
                    if bg2 or border2:
                        items.append(
                            Item(
                                kind="shape",
                                x=x,
                                y=y,
                                w=w,
                                h=h,
                                rotation=rot,
                                z=z,
                                bg=bg2,
                                border=border2,
                                border_w=border_w2,
                            )
                        )

        slides_items.append(items)

    # HTML generation
    def item_style(it: Item) -> str:
        # Position in percentages of the PPT base size.
        x_pct = (it.x / slide_w) * 100.0
        y_pct = (it.y / slide_h) * 100.0
        w_pct = (it.w / slide_w) * 100.0
        h_pct = (it.h / slide_h) * 100.0

        parts = [
            f"--x:{x_pct:.6f}%;",
            f"--y:{y_pct:.6f}%;",
            f"--w:{w_pct:.6f}%;",
            f"--h:{h_pct:.6f}%;",
            f"z-index:{it.z};",
        ]
        if it.font_px is not None:
            parts.append(f"--fs:{it.font_px:.3f}px;")
        # Keep PPT layout faithful: only add padding when there's a visible box background/border.
        if it.kind == "text":
            pad = "0.32em" if (it.bg or it.border) else "0"
            parts.append(f"--pad:{pad};")
        if it.bg:
            parts.append(f"--bg:{it.bg};")
        if it.border:
            parts.append(f"--bd:{it.border};")
        if it.border_w is not None:
            parts.append(f"--bw:{it.border_w:.2f}px;")
        if it.color:
            parts.append(f"--tc:{it.color};")
        if it.rotation:
            parts.append(f"--rot:{it.rotation:.4f}deg;")
        return "".join(parts)

    slides_html: List[str] = []
    for idx, items in enumerate(slides_items, start=1):
        node_items: List[str] = []
        for it in items:
            if it.kind == "image":
                node_items.append(
                    f"""
          <button class="ppt-item ppt-img" style="{item_style(it)}" type="button" data-lightbox="{html.escape(it.src or '')}" aria-label="Open image">
            <img src="{html.escape(it.src or '')}" alt="{html.escape(it.alt)}" loading="lazy" />
          </button>
                    """.rstrip()
                )
            elif it.kind == "text":
                weight = 700 if it.bold else 500
                node_items.append(
                    f"""
          <div class="ppt-item ppt-text" style="{item_style(it)}">
            <div class="ppt-text-inner" style="font-weight:{weight};">{it.text}</div>
          </div>
                    """.rstrip()
                )
            elif it.kind == "shape":
                node_items.append(
                    f"""<div class="ppt-item ppt-shape" style="{item_style(it)}" aria-hidden="true"></div>"""
                )
            elif it.kind == "line":
                node_items.append(
                    f"""<div class="ppt-item ppt-line" style="{item_style(it)}" aria-hidden="true"></div>"""
                )

        slides_html.append(
            f"""
    <section class="slide" id="slide-{idx}" data-label="P{idx:02d}">
      <div class="slide-content">
        <div class="artboard" data-base-w="{base_w_px:.3f}" data-base-h="{base_h_px:.3f}" style="--ratio:{ratio:.6f};">
{os.linesep.join(node_items)}
        </div>
      </div>
    </section>
            """.rstrip()
        )

    out_html.write_text(
        f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>智能体平台 QMoss + QDify（PPT → HTML）</title>

  <!-- Fonts (distinctive, but the PPT layout drives typography) -->
  <link rel="stylesheet" href="https://api.fontshare.com/v2/css?f[]=satoshi@400,500,700&f[]=jet-brains-mono@400,600&display=swap" />

  <style>
    /* ===========================================
       VIEWPORT FITTING: MANDATORY BASE STYLES
       =========================================== */
    html, body {{
      height: 100%;
      overflow-x: hidden;
    }}

    html {{
      scroll-snap-type: y mandatory;
      scroll-behavior: smooth;
    }}

    :root {{
      --bg-0: #070911;
      --bg-1: #0b1020;
      --glass: rgba(255, 255, 255, 0.06);
      --glass-2: rgba(255, 255, 255, 0.03);
      --line: rgba(255, 255, 255, 0.12);
      --text-0: rgba(255, 255, 255, 0.92);
      --text-1: rgba(255, 255, 255, 0.70);
      --accent: #3dffcf;

      --font-body: "Satoshi", ui-sans-serif, system-ui, -apple-system, "Segoe UI", sans-serif;
      --font-mono: "JetBrains Mono", ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, "Liberation Mono", monospace;

      --slide-padding: clamp(0.75rem, 3.4vw, 3rem);
      --content-gap: clamp(0.6rem, 1.8vw, 1.3rem);
      --ui-size: clamp(0.68rem, 1vw, 0.92rem);
    }}

    * {{ margin: 0; padding: 0; box-sizing: border-box; }}

    body {{
      font-family: var(--font-body);
      background: radial-gradient(1200px 900px at 15% 10%, rgba(124, 92, 255, 0.18), transparent 60%),
                  radial-gradient(900px 700px at 80% 20%, rgba(61, 255, 207, 0.12), transparent 55%),
                  radial-gradient(900px 700px at 65% 88%, rgba(255, 79, 216, 0.10), transparent 55%),
                  linear-gradient(180deg, var(--bg-0), var(--bg-1));
      color: var(--text-0);
      height: 100%;
      overflow-x: hidden;
    }}

    .slide {{
      width: 100vw;
      height: 100vh;
      height: 100dvh;
      overflow: hidden;
      scroll-snap-align: start;
      display: flex;
      flex-direction: column;
      position: relative;
      padding: var(--slide-padding);
    }}

    .slide-content {{
      flex: 1;
      display: grid;
      place-items: center;
      max-height: 100%;
      overflow: hidden;
    }}

    /* PPT artboard: preserve 16:9-ish layout and place elements in percentages. */
    .artboard {{
      position: relative;
      width: min(96vw, calc(92vh * var(--ratio)));
      aspect-ratio: var(--ratio);
      background: linear-gradient(180deg, var(--glass), var(--glass-2));
      border: 1px solid var(--line);
      border-radius: 18px;
      box-shadow: 0 30px 140px rgba(0, 0, 0, 0.55);
      overflow: hidden;
      transform: translateZ(0);
    }}

    /* Items are positioned via CSS vars in %; font sizes scale with artboard height. */
    .ppt-item {{
      position: absolute;
      left: var(--x);
      top: var(--y);
      width: var(--w);
      height: var(--h);
      transform: rotate(var(--rot, 0deg));
      transform-origin: center;
    }}

    .ppt-text {{
      display: block;
      color: var(--tc, rgba(255,255,255,0.92));
      background: var(--bg, transparent);
      border: var(--bw, 0px) solid var(--bd, transparent);
      border-radius: 12px;
      overflow: hidden;
      padding: var(--pad, 0);
      backdrop-filter: blur(8px);
      -webkit-backdrop-filter: blur(8px);
    }}

    .ppt-text-inner {{
      width: 100%;
      height: 100%;
      font-size: calc(var(--s, 1) * var(--fs, 18px));
      line-height: 1.15;
      color: inherit;
      overflow: hidden;
      word-break: break-word;
      white-space: normal;
    }}

    .ppt-shape {{
      background: var(--bg, transparent);
      border: var(--bw, 1px) solid var(--bd, rgba(255,255,255,0.15));
      border-radius: 12px;
      opacity: 0.9;
    }}

    .ppt-line {{
      background: transparent;
      border-top: var(--bw, 2px) solid var(--bd, rgba(255,255,255,0.55));
      height: 0;
      width: var(--w);
    }}

    /* Clickable images with full-size lightbox */
    .ppt-img {{
      padding: 0;
      border: none;
      background: transparent;
      cursor: zoom-in;
    }}
    .ppt-img img {{
      width: 100%;
      height: 100%;
      object-fit: contain;
      display: block;
      border-radius: 10px;
      background: rgba(0, 0, 0, 0.18);
      box-shadow: 0 18px 60px rgba(0, 0, 0, 0.40);
      transition: transform 180ms ease;
    }}
    .ppt-img:hover img {{
      transform: scale(1.01);
    }}

    /* UI */
    .progress {{
      position: fixed;
      top: 0;
      left: 0;
      height: 3px;
      width: 100vw;
      z-index: 50;
      background: rgba(255,255,255,0.08);
      backdrop-filter: blur(8px);
      -webkit-backdrop-filter: blur(8px);
    }}
    .progress > div {{
      height: 100%;
      width: 0%;
      background: linear-gradient(90deg, #3dffcf, #ff4fd8, #7c5cff);
      box-shadow: 0 0 22px rgba(61, 255, 207, 0.26);
      transition: width 220ms ease;
    }}

    .nav-dots {{
      position: fixed;
      right: clamp(0.55rem, 1.6vw, 1.2rem);
      top: 50%;
      transform: translateY(-50%);
      z-index: 50;
      display: flex;
      flex-direction: column;
      gap: 9px;
      padding: 10px 8px;
      border-radius: 999px;
      background: rgba(0,0,0,0.18);
      border: 1px solid rgba(255,255,255,0.10);
      backdrop-filter: blur(12px);
      -webkit-backdrop-filter: blur(12px);
    }}
    .nav-dots button {{
      width: 10px;
      height: 10px;
      border-radius: 999px;
      border: 1px solid rgba(255,255,255,0.20);
      background: rgba(255,255,255,0.08);
      cursor: pointer;
      transition: transform 220ms ease, background 220ms ease, border-color 220ms ease;
    }}
    .nav-dots button[aria-current="true"] {{
      background: linear-gradient(135deg, #3dffcf, #ff4fd8);
      border-color: rgba(61,255,207,0.55);
      transform: scale(1.25);
      box-shadow: 0 0 18px rgba(61,255,207,0.22);
    }}

    .counter {{
      position: fixed;
      left: clamp(0.65rem, 2vw, 1.4rem);
      bottom: clamp(0.65rem, 2vw, 1.2rem);
      z-index: 50;
      font-family: var(--font-mono);
      font-size: var(--ui-size);
      color: rgba(255,255,255,0.66);
      background: rgba(0,0,0,0.18);
      border: 1px solid rgba(255,255,255,0.10);
      border-radius: 999px;
      padding: 0.5em 0.75em;
      backdrop-filter: blur(10px);
      -webkit-backdrop-filter: blur(10px);
    }}

    /* Lightbox */
    .lightbox {{
      position: fixed;
      inset: 0;
      z-index: 200;
      display: none;
      place-items: center;
      background: rgba(0,0,0,0.74);
      backdrop-filter: blur(10px);
      -webkit-backdrop-filter: blur(10px);
      padding: clamp(0.8rem, 2.5vw, 2rem);
    }}
    .lightbox.open {{
      display: grid;
    }}
    .lightbox-card {{
      width: min(96vw, 1400px);
      height: min(92vh, 900px);
      border-radius: 16px;
      border: 1px solid rgba(255,255,255,0.12);
      background: rgba(10, 12, 20, 0.55);
      box-shadow: 0 40px 200px rgba(0,0,0,0.65);
      overflow: hidden;
      display: grid;
      grid-template-rows: auto 1fr;
    }}
    .lightbox-bar {{
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 12px;
      padding: 10px 12px;
      border-bottom: 1px solid rgba(255,255,255,0.10);
    }}
    .lightbox-bar .title {{
      font-family: var(--font-mono);
      font-size: var(--ui-size);
      color: rgba(255,255,255,0.70);
      overflow: hidden;
      text-overflow: ellipsis;
      white-space: nowrap;
      max-width: 70%;
    }}
    .lightbox-actions {{
      display: inline-flex;
      gap: 10px;
      align-items: center;
    }}
    .lightbox-actions a,
    .lightbox-actions button {{
      font-family: var(--font-mono);
      font-size: var(--ui-size);
      color: rgba(255,255,255,0.82);
      background: rgba(255,255,255,0.06);
      border: 1px solid rgba(255,255,255,0.12);
      border-radius: 999px;
      padding: 0.45em 0.75em;
      cursor: pointer;
    }}
    .lightbox-actions a:hover,
    .lightbox-actions button:hover {{
      border-color: rgba(61,255,207,0.32);
    }}
    .lightbox-stage {{
      display: grid;
      place-items: center;
      padding: 12px;
      overflow: hidden;
    }}
    .lightbox-stage img {{
      width: 100%;
      height: 100%;
      object-fit: contain;
      background: rgba(0,0,0,0.25);
      border-radius: 12px;
    }}

    /* ===========================================
       RESPONSIVE BREAKPOINTS (mandatory)
       =========================================== */
    @media (max-height: 700px) {{
      :root {{
        --slide-padding: clamp(0.55rem, 2.6vw, 1.6rem);
      }}
    }}
    @media (max-height: 600px) {{
      .nav-dots {{ display: none; }}
    }}
    @media (max-height: 500px) {{
      .progress {{ display: none; }}
    }}
    @media (max-width: 900px) {{
      .nav-dots {{ right: 0.55rem; }}
      .counter {{ left: 0.65rem; }}
      .artboard {{ width: min(96vw, calc(88vh * var(--ratio))); }}
    }}

    @media (prefers-reduced-motion: reduce) {{
      html {{ scroll-behavior: auto; }}
      * {{ transition-duration: 0.01ms !important; animation-duration: 0.01ms !important; }}
    }}
  </style>
</head>
<body>
  <div class="progress" aria-hidden="true"><div id="progressFill"></div></div>
  <nav class="nav-dots" aria-label="Slide navigation" id="navDots"></nav>
  <div class="counter" id="counter" aria-live="polite">01 / {len(slides_items):02d}</div>

  <!-- Slides -->
{os.linesep.join(slides_html)}

  <!-- Lightbox -->
  <div class="lightbox" id="lightbox" aria-hidden="true">
    <div class="lightbox-card" role="dialog" aria-modal="true" aria-label="Image viewer">
      <div class="lightbox-bar">
        <div class="title" id="lightboxTitle">image</div>
        <div class="lightbox-actions">
          <a id="lightboxOpen" href="#" target="_blank" rel="noreferrer">Open</a>
          <button type="button" id="lightboxClose">Close (Esc)</button>
        </div>
      </div>
      <div class="lightbox-stage" id="lightboxStage">
        <img id="lightboxImg" alt="Full-size preview" />
      </div>
    </div>
  </div>

  <script>
    class SlidePresentation {{
      constructor() {{
        this.slides = Array.from(document.querySelectorAll(".slide"));
        this.navDots = document.getElementById("navDots");
        this.progressFill = document.getElementById("progressFill");
        this.counter = document.getElementById("counter");
        this.reduceMotion = window.matchMedia && window.matchMedia("(prefers-reduced-motion: reduce)").matches;
        this.activeIndex = 0;
        this._wheelLock = false;

        this._buildDots();
        this._bind();
        this._syncActiveFromScroll();
        this._installArtboardScaling();
      }}

      _installArtboardScaling() {{
        const boards = Array.from(document.querySelectorAll(".artboard"));
        const updateOne = (b) => {{
          const baseH = parseFloat(b.getAttribute("data-base-h") || "0");
          if (!baseH) return;
          const scale = b.clientHeight / baseH;
          b.style.setProperty("--s", String(scale));
        }};
        const updateAll = () => boards.forEach(updateOne);
        updateAll();
        window.addEventListener("resize", () => updateAll(), {{ passive: true }});
      }}

      _buildDots() {{
        this.navDots.innerHTML = "";
        this.slides.forEach((slide, index) => {{
          const label = slide.getAttribute("data-label") || `Slide ${{index + 1}}`;
          const button = document.createElement("button");
          button.type = "button";
          button.setAttribute("aria-label", `Go to ${{label}}`);
          button.addEventListener("click", () => this.goTo(index));
          this.navDots.appendChild(button);
        }});
      }}

      _bind() {{
        window.addEventListener("keydown", (e) => {{
          const key = e.key;
          const isNavKey = ["ArrowDown","ArrowRight","PageDown"," ","Enter","ArrowUp","ArrowLeft","PageUp","Home","End"].includes(key);
          if (!isNavKey) return;
          e.preventDefault();
          if (key === "Home") return this.goTo(0);
          if (key === "End") return this.goTo(this.slides.length - 1);
          if (key === "ArrowUp" || key === "ArrowLeft" || key === "PageUp") return this.prev();
          return this.next();
        }}, {{ passive: false }} );

        window.addEventListener("wheel", (e) => {{
          if (this._wheelLock) return;
          const dy = e.deltaY || 0;
          if (Math.abs(dy) < 14) return;
          this._wheelLock = true;
          if (dy > 0) this.next(); else this.prev();
          window.setTimeout(() => {{ this._wheelLock = false; }}, 360);
        }}, {{ passive: true }});

        window.addEventListener("scroll", () => this._syncActiveFromScroll(), {{ passive: true }});
      }}

      _syncActiveFromScroll() {{
        const viewportMid = window.scrollY + window.innerHeight / 2;
        let bestIndex = 0;
        let bestDist = Infinity;
        for (let i = 0; i < this.slides.length; i++) {{
          const rect = this.slides[i].getBoundingClientRect();
          const slideMid = window.scrollY + rect.top + rect.height / 2;
          const dist = Math.abs(slideMid - viewportMid);
          if (dist < bestDist) {{ bestDist = dist; bestIndex = i; }}
        }}
        this.activeIndex = bestIndex;
        this._renderUI();
      }}

      _renderUI() {{
        const total = this.slides.length;
        const current = this.activeIndex + 1;
        const pad2 = (n) => String(n).padStart(2, "0");
        this.counter.textContent = `${{pad2(current)}} / ${{pad2(total)}}`;
        const percent = (this.activeIndex / Math.max(1, total - 1)) * 100;
        this.progressFill.style.width = `${{percent}}%`;
        const dots = Array.from(this.navDots.querySelectorAll("button"));
        dots.forEach((b, i) => b.setAttribute("aria-current", i === this.activeIndex ? "true" : "false"));
      }}

      goTo(index) {{
        const slide = this.slides[index];
        if (!slide) return;
        slide.scrollIntoView({{ behavior: this.reduceMotion ? "auto" : "smooth", block: "start" }});
        this.activeIndex = Math.max(0, Math.min(index, this.slides.length - 1));
        this._renderUI();
      }}

      next() {{ this.goTo(this.activeIndex + 1); }}
      prev() {{ this.goTo(this.activeIndex - 1); }}
    }}

    class Lightbox {{
      constructor() {{
        this.el = document.getElementById("lightbox");
        this.img = document.getElementById("lightboxImg");
        this.title = document.getElementById("lightboxTitle");
        this.openLink = document.getElementById("lightboxOpen");
        this.closeBtn = document.getElementById("lightboxClose");

        this.closeBtn.addEventListener("click", () => this.close());
        this.el.addEventListener("click", (e) => {{
          if (e.target === this.el) this.close();
        }});
        window.addEventListener("keydown", (e) => {{
          if (e.key === "Escape") this.close();
        }});

        document.body.addEventListener("click", (e) => {{
          const btn = e.target && e.target.closest && e.target.closest("[data-lightbox]");
          if (!btn) return;
          const src = btn.getAttribute("data-lightbox");
          if (!src) return;
          this.open(src);
        }});
      }}

      open(src) {{
        this.img.src = src;
        this.title.textContent = src.split("/").pop() || "image";
        this.openLink.href = src;
        this.el.classList.add("open");
        this.el.setAttribute("aria-hidden", "false");
      }}

      close() {{
        if (!this.el.classList.contains("open")) return;
        this.el.classList.remove("open");
        this.el.setAttribute("aria-hidden", "true");
        this.img.src = "";
      }}
    }}

    new SlidePresentation();
    new Lightbox();
  </script>
</body>
</html>
""",
        encoding="utf-8",
    )


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    pptx_path = root / "对外分享" / "智能体平台QMoss+QDify.pptx"
    out_html = root / "对外分享" / "智能体平台QMoss+QDify.html"
    assets_dir = root / "对外分享" / "智能体平台QMoss+QDify-assets"
    build_deck(pptx_path, out_html, assets_dir)
    print(f"OK: wrote {out_html}")


if __name__ == "__main__":
    main()
